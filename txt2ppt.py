#-*- coding: utf-8 -*-
import fix_path # has to be first.

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor

import lxml
from lxml import etree

import re
from __builtin__ import True

class Slides:
    prs = Presentation()
    file_name = "slides.pptx"
    blank_slide_layout = prs.slide_layouts[6]
    font_size = Pt(30)
    font_color = "000000"
    background_color = "FFFFFF"

    def __init__(self, file_name):
        self.prs = Presentation()
        self.file_name = file_name    

    def set_font_size(self, size):
        self.font_size = Pt(size)

    def set_font_color(self, color):
        self.font_color = color
        
    def set_background_color(self, color):
        self.background_color = color
        
    def new(self, text):
        slide = self.prs.slides.add_slide(self.blank_slide_layout)
        
        left = Inches (0.5)
        top = Inches (2)
        width = Inches (9)
        height = Inches (3)
    
        txBox = slide.shapes.add_textbox(left, top, width, height)
        #Set background color
        self.slide_bg_color(slide, self.background_color)
        
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        p = tf.add_paragraph()
        p.font.size = self.font_size
        p.font.color.rgb = RGBColor.from_string(self.font_color)
        p.alignment = PP_ALIGN.CENTER
        
        for part in self.get_text_parts(text):
            run = p.add_run()
            run.text = part['text']
            if part['type'] == 'bold':
                run.font.bold = True
            if part['type'] == 'italic':
                run.font.italic = True
        
    def slide_bg_color(self, slide, bgcolor):
        """
        Ideal é colocar dentro da lib, no objeto slide. Mas preferi colocar aqui para não necessitar manter controle da versão da lib.
        Dessa forma, a lib pptx continua a mesma
        """
        #<p:bg><p:bgPr><a:solidFill><a:srgbClr val="010101"/></a:solidFill><a:effectLst/></p:bgPr></p:bg>
        el_bg = etree.Element('{http://schemas.openxmlformats.org/presentationml/2006/main}bg', nsmap={'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
        slide._element[0].insert(0, el_bg)
        el_bg_pr = etree.Element('{http://schemas.openxmlformats.org/presentationml/2006/main}bgPr', nsmap={'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
        el_bg.append(el_bg_pr)
        el_solid_fill = etree.Element('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill', nsmap={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
        el_bg_pr.append(el_solid_fill)
        color = etree.Element('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr', val=bgcolor, nsmap={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
        el_solid_fill.append(color)
        
    def get_text_parts(self, text):
        """
        Return a list of dict with parts of text with information about bold and italic
        """
        #First check if there a satisfied number of *. If not, raise exception
        cont = text.count('*')
        if cont % 2 != 0:
            raise Exception('Invalid syntax for line: ' + text)
        
        ret = []
        p = re.compile('\*\*.+?\*\*|\*.+?\*')
        bold = re.compile('\*\*.+?\*\*')
        current_index = 0
        for line in p.finditer(text):
            start = line.start()
            end = start + len(line.group())
            #Pega o normal imediatamente anterior
            if start > current_index:
                ret.append({'text': text[current_index:start], 'type': 'normal'})
            current_index = end
            if bold.match(line.group()):
                type = 'bold'
            else:
                type = 'italic'
            ret.append({'text': line.group().replace('*', ''), 'type': type})
        
        if current_index < len(text):
            ret.append({'text': text[current_index:len(text)], 'type': 'normal'})
        return ret
        
    def save(self):
        self.prs.save(self.file_name)